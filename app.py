from flask import Flask, render_template, request, session, redirect, url_for
import json, os, fitz
from pptx import Presentation
from groq import Groq
from dotenv import load_dotenv

#for importing libraries for creating pdfs
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch 
import uuid 
from flask import send_from_directory

PDF_FOLDER = "generated_pdfs"
os.makedirs(PDF_FOLDER, exist_ok=True)

def create_questions_pdf(questions, filepath):
    doc = SimpleDocTemplate(filepath, pagesize=letter,
                             topMargin=0.7*inch, bottomMargin=0.7*inch)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("Practice Questions", styles['Title']))
    story.append(Spacer(1, 16))

    for q in questions:
        story.append(Paragraph(f"{q['id']}. {q['question']}", styles['Heading3']))
        story.append(Spacer(1, 4))

        for i, option in enumerate(q['options']):
            letter_label = chr(65 + i)  # A, B, C, D
            story.append(Paragraph(f"{letter_label}. {option}", styles['Normal']))

        story.append(Spacer(1, 6))
        story.append(Paragraph(f"<b>Answer:</b> {q['answer']}", styles['Normal']))
        story.append(Spacer(1, 16))

    doc.build(story)

app = Flask(__name__)
app.secret_key = "quiz_secret_key"
UPLOAD_FOLDER = "uploads"

load_dotenv()
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

client = Groq(api_key=os.getenv("GROQ_API_KEY"))


# Extract the text from uploaded file 
def extract_text(filepath):
    if filepath.endswith(".pdf"):
        doc = fitz.open(filepath)
        return "\n".join(page.get_text() for page in doc)

    elif filepath.endswith(".pptx"):
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text_frame.text)
        return "\n".join(text)

    return ""

#no_of_questions = ""

def generate_questions(text, no_of_questions):
    prompt = f"""
Generate {no_of_questions} multiple choice questions from this lecture content.
Return ONLY a valid JSON array, no explanation, no markdown, no backticks.
Each object must have exactly these fields:
- "id": number
- "question": string
- "options": array of 4 strings
- "answer": string (must exactly match one of the options)

Lecture content:
{text[:8000]}
"""
    response = client.chat.completions.create(
        model="openai/gpt-oss-120b",
        messages=[{"role": "user", "content": prompt}]
    )
    raw = response.choices[0].message.content.strip()

    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1]
    if raw.endswith("```"):
        raw = raw.rsplit("```", 1)[0]

    return json.loads(raw.strip())

# ── Routes ────────────────────────────────────────────────────
@app.route("/", methods=["GET", "POST"])
def upload():
    if request.method == "GET":
        session.clear()

    if request.method == "POST":

        session.clear() #clear session to clear previous data

        no_of_questions = request.form.get("num_questions")

        file = request.files.get("file")

        if not file or file.filename == "":
            return render_template("upload.html", error="Please select a file.")

        filename = file.filename
        if not (filename.endswith(".pdf") or filename.endswith(".pptx")):
            return render_template("upload.html", error="Only PDF and PPTX files are supported.")

        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        try:
            text = extract_text(filepath)
            if not text.strip():
                return render_template("upload.html", error="Could not extract text from file. Try another file.")

            questions = generate_questions(text, no_of_questions)

            session["questions"] = questions #adding sessions so there will be differnt users even without login

            # Generate PDF right away
            pdf_filename = f"questions_{uuid.uuid4().hex}.pdf"
            pdf_path = os.path.join(PDF_FOLDER, pdf_filename)
            create_questions_pdf(questions, pdf_path)
            session["pdf_filename"] = pdf_filename
            
            return redirect(url_for("start"))

        except Exception as e:
            return render_template("upload.html", error=f"Something went wrong: {str(e)}")

    return render_template("upload.html")

@app.route("/generated")
def start():
    session["current"] = 0
    session["score"] = 0
    session["answers"] = []
    questions = session.get("questions") #get the stored questions from session
    if not questions:
        return render_template("upload.html")

    index = session.get("current", 0)
    question = questions[index]
    total = len(questions)
    print("Total questions:", total)  # Debugging line to check the total number of questions

    return render_template("quiz-generated.html", total=total, question=question, index=index)

@app.route("/quiz", methods=["GET", "POST"])
def quiz():
    questions = session.get("questions") #get the stored questions from session
    if not questions:
        return render_template("upload.html")

    index = session.get("current", 0)

    if request.method == "POST":
        selected = request.form.get("answer")
        correct = questions[index]["answer"]

        session["answers"] = session.get("answers", []) + [selected]

        if selected == correct:
            session["score"] += 1

        session["current"] += 1
        index = session["current"]

        if index >= len(questions):
            return redirect(url_for("result"))

    question = questions[index]
    total = len(questions)
    return render_template("quiz.html", question=question, index=index, total=total)

@app.route("/result", methods=["GET", "POST"])
def result():
    score = session.get("score", 0)
    questions = session.get("questions") #get the stored questions from session
    total = len(questions)
    return render_template("result.html", score=score, total=total)

@app.route("/home")
def home():
    return render_template("upload.html")

def load_questions():
    with open("questions.json") as f:
        return json.load(f)

#endpoint for downloading the generated PDF
@app.route("/download")
def download_pdf():
    pdf_filename = session.get("pdf_filename")
    if not pdf_filename:
        return "No file available.", 404
    return send_from_directory(PDF_FOLDER, pdf_filename,
                                as_attachment=True,
                                download_name="practice_questions.pdf")

if __name__ == "__main__":
    app.run(debug=True)